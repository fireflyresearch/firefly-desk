# Copyright 2026 Firefly Software Solutions Inc
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0

"""File upload REST API."""

from __future__ import annotations

import uuid
from typing import Annotated

from fastapi import APIRouter, Depends, Form, HTTPException, Query, Request, Response, UploadFile
from fastapi.responses import Response as FastAPIResponse

from flydesk.files.extractor import ContentExtractor
from flydesk.files.models import FileUpload
from flydesk.files.repository import FileUploadRepository
from flydesk.files.storage import FileStorageProvider

router = APIRouter(tags=["files"])


# ---------------------------------------------------------------------------
# Dependencies
# ---------------------------------------------------------------------------


def get_file_repo() -> FileUploadRepository:
    """Provide a FileUploadRepository instance.

    In production this is wired to the real database session factory.
    In tests the dependency is overridden via app.dependency_overrides.
    """
    raise NotImplementedError(
        "get_file_repo must be overridden via app.dependency_overrides"
    )


def get_file_storage() -> FileStorageProvider:
    """Provide a FileStorageProvider instance."""
    raise NotImplementedError(
        "get_file_storage must be overridden via app.dependency_overrides"
    )


def get_content_extractor() -> ContentExtractor:
    """Provide a ContentExtractor instance."""
    raise NotImplementedError(
        "get_content_extractor must be overridden via app.dependency_overrides"
    )


FileRepo = Annotated[FileUploadRepository, Depends(get_file_repo)]
Storage = Annotated[FileStorageProvider, Depends(get_file_storage)]
Extractor = Annotated[ContentExtractor, Depends(get_content_extractor)]


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------


@router.post("/api/chat/upload", status_code=201)
async def upload_file(
    request: Request,
    file: UploadFile,
    repo: FileRepo,
    storage: Storage,
    extractor: Extractor,
    conversation_id: Annotated[str | None, Form()] = None,
) -> dict:
    """Upload a file (multipart/form-data).

    Accepts a ``file`` and an optional ``conversation_id`` form field.
    """
    config = getattr(request.app.state, "config", None)
    max_size = (config.file_max_size_mb if config else 50) * 1024 * 1024

    # Read file content
    content = await file.read()

    # Validate file size
    if len(content) > max_size:
        raise HTTPException(
            status_code=413,
            detail=f"File exceeds maximum size of {config.file_max_size_mb if config else 50}MB",
        )

    filename = file.filename or "unnamed"
    content_type = file.content_type or "application/octet-stream"

    # Store the file
    storage_path = await storage.store(filename, content, content_type)

    # Extract text content
    extracted_text = await extractor.extract(filename, content, content_type)

    # Get user info
    user_session = getattr(request.state, "user_session", None)
    user_id = user_session.user_id if user_session else "anonymous"

    # Create the upload record
    file_upload = FileUpload(
        id=str(uuid.uuid4()),
        conversation_id=conversation_id,
        user_id=user_id,
        filename=filename,
        content_type=content_type,
        file_size=len(content),
        storage_path=storage_path,
        storage_backend="local",
        extracted_text=extracted_text,
    )
    await repo.create(file_upload)

    return {
        "id": file_upload.id,
        "filename": file_upload.filename,
        "content_type": file_upload.content_type,
        "file_size": file_upload.file_size,
        "conversation_id": file_upload.conversation_id,
        "extracted_text": file_upload.extracted_text,
    }


@router.get("/api/files/{file_id}")
async def get_file_metadata(file_id: str, repo: FileRepo) -> dict:
    """Return metadata for a single file upload."""
    upload = await repo.get(file_id)
    if upload is None:
        raise HTTPException(status_code=404, detail=f"File {file_id} not found")
    return {
        "id": upload.id,
        "filename": upload.filename,
        "content_type": upload.content_type,
        "file_size": upload.file_size,
        "conversation_id": upload.conversation_id,
        "storage_backend": upload.storage_backend,
        "extracted_text": upload.extracted_text,
        "metadata": upload.metadata,
        "created_at": upload.created_at.isoformat() if upload.created_at else None,
    }


@router.get("/api/files/{file_id}/download")
async def download_file(file_id: str, repo: FileRepo, storage: Storage) -> FastAPIResponse:
    """Download a file's content."""
    upload = await repo.get(file_id)
    if upload is None:
        raise HTTPException(status_code=404, detail=f"File {file_id} not found")

    content = await storage.retrieve(upload.storage_path)
    return Response(
        content=content,
        media_type=upload.content_type,
        headers={
            "Content-Disposition": f'attachment; filename="{upload.filename}"',
        },
    )


@router.delete("/api/files/{file_id}", status_code=204)
async def delete_file(file_id: str, repo: FileRepo, storage: Storage) -> Response:
    """Delete a file and its storage."""
    upload = await repo.get(file_id)
    if upload is None:
        raise HTTPException(status_code=404, detail=f"File {file_id} not found")

    await storage.delete(upload.storage_path)
    await repo.delete(file_id)
    return Response(status_code=204)


@router.get("/api/files/{file_id}/render")
async def render_file(
    file_id: str,
    repo: FileRepo,
    storage: Storage,
    render_format: str = Query(
        default="html", alias="format", description="Render format: html, json, text"
    ),
) -> dict:
    """Render a file for preview in the file viewer widget."""
    upload = await repo.get(file_id)
    if upload is None:
        raise HTTPException(status_code=404, detail="File not found")

    content = await storage.retrieve(upload.storage_path)
    ct = upload.content_type

    # Route by content type
    if ct in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    }:
        return _render_docx(content)
    if ct in {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    }:
        return _render_xlsx(content)
    if ct in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }:
        return _render_pptx(content)
    if ct == "application/pdf":
        return {"type": "pdf", "download_url": f"/api/files/{file_id}/download"}
    if ct.startswith("image/"):
        return {"type": "image", "download_url": f"/api/files/{file_id}/download"}

    # Fallback: text
    text = content.decode("utf-8", errors="replace")
    return {"type": "text", "text": text}


# ---------------------------------------------------------------------------
# Render helpers
# ---------------------------------------------------------------------------


def _render_docx(content: bytes) -> dict:
    """Convert DOCX content to an HTML representation."""
    import io

    from docx import Document

    doc = Document(io.BytesIO(content))
    html_parts: list[str] = []
    for para in doc.paragraphs:
        if para.style and para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading ", "").strip()
            if level.isdigit():
                html_parts.append(f"<h{level}>{para.text}</h{level}>")
            else:
                html_parts.append(f"<p><strong>{para.text}</strong></p>")
        elif para.text.strip():
            html_parts.append(f"<p>{para.text}</p>")

    for table in doc.tables:
        rows_html: list[str] = []
        for i, row in enumerate(table.rows):
            tag = "th" if i == 0 else "td"
            cells = "".join(f"<{tag}>{cell.text}</{tag}>" for cell in row.cells)
            rows_html.append(f"<tr>{cells}</tr>")
        html_parts.append(f"<table>{''.join(rows_html)}</table>")

    return {"type": "docx", "html": "\n".join(html_parts)}


def _render_xlsx(content: bytes) -> dict:
    """Convert XLSX content to a structured JSON table representation."""
    import io

    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows_data = list(ws.iter_rows(values_only=True))
        if not rows_data:
            continue
        header = rows_data[0]
        columns = [
            {"key": f"col_{i}", "label": str(h) if h is not None else f"Column {i + 1}"}
            for i, h in enumerate(header)
        ]
        rows = []
        for row in rows_data[1:]:
            row_dict = {}
            for i, cell in enumerate(row):
                row_dict[f"col_{i}"] = cell if cell is not None else ""
            rows.append(row_dict)
        sheets.append({"title": ws.title, "columns": columns, "rows": rows})
    wb.close()
    return {"type": "xlsx", "sheets": sheets}


def _render_pptx(content: bytes) -> dict:
    """Convert PPTX content to a slide-by-slide text representation."""
    import io

    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        slides.append({"slide_number": i, "texts": texts})
    return {"type": "pptx", "slides": slides}
