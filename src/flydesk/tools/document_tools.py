# Copyright 2026 Firefly Software Solutions Inc
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0

"""Document tools for reading, creating, modifying, and converting office documents.

These tools let the agent work with PDF, DOCX, XLSX, and PPTX files
stored in the platform's file storage backend.
"""

from __future__ import annotations

import csv
import io
import json
import logging
import uuid
from typing import TYPE_CHECKING, Any

from flydesk.catalog.enums import HttpMethod, RiskLevel
from flydesk.tools.builtin import BUILTIN_SYSTEM_ID
from flydesk.tools.factory import ToolDefinition

from flydesk.settings.models import LLMRuntimeSettings

if TYPE_CHECKING:
    from flydesk.files.repository import FileUploadRepository
    from flydesk.files.storage import FileStorageProvider
    from flydesk.settings.repository import SettingsRepository

logger = logging.getLogger(__name__)


def _coerce_dict(value: Any) -> dict[str, Any]:
    """Coerce *value* to a dict.

    LLMs sometimes send JSON-object parameters as a serialised string rather
    than an actual dict.  This helper transparently handles both cases.
    """
    if isinstance(value, dict):
        return value
    if isinstance(value, str):
        try:
            parsed = json.loads(value)
            if isinstance(parsed, dict):
                return parsed
        except (json.JSONDecodeError, TypeError):
            logger.debug("Could not parse string value as JSON dict")
    return {}


# ---------------------------------------------------------------------------
# Tool definitions
# ---------------------------------------------------------------------------


def document_read_tool() -> ToolDefinition:
    """Tool definition for reading document content."""
    return ToolDefinition(
        endpoint_id="__builtin__document_read",
        name="document_read",
        description=(
            "Extract text, tables, and metadata from a document file "
            "(PDF, DOCX, XLSX, PPTX). Returns the textual content of the document. "
            "For large documents, use page_start/page_end to read specific page "
            "ranges (PDF) or set max_chars to limit output. If the output is "
            "truncated, call again with the next page range to continue reading. "
            "Prefer using file_id (when available) over file_path for "
            "unambiguous file resolution."
        ),
        risk_level=RiskLevel.READ,
        system_id=BUILTIN_SYSTEM_ID,
        method=HttpMethod.GET.value,
        path="/__builtin__/documents/read",
        parameters={
            "file_id": {
                "type": "string",
                "description": (
                    "Unique file upload ID. Use this when provided in file "
                    "context hints. Takes priority over file_path."
                ),
                "required": False,
            },
            "file_path": {
                "type": "string",
                "description": (
                    "Storage path or filename of the document to read. "
                    "Used as fallback when file_id is not available."
                ),
                "required": False,
            },
            "page_start": {
                "type": "integer",
                "description": "First page to read (1-based, PDF only). Defaults to 1.",
                "required": False,
            },
            "page_end": {
                "type": "integer",
                "description": "Last page to read (inclusive, PDF only). Defaults to all pages.",
                "required": False,
            },
            "max_chars": {
                "type": "integer",
                "description": "Maximum characters to return. Defaults to 30000.",
                "required": False,
            },
        },
    )


# Fallback max chars for document_read; matches LLMRuntimeSettings.document_read_max_chars default.
_DEFAULT_MAX_CHARS = LLMRuntimeSettings().document_read_max_chars


def document_create_tool() -> ToolDefinition:
    """Tool definition for creating new documents."""
    return ToolDefinition(
        endpoint_id="__builtin__document_create",
        name="document_create",
        description=(
            "Generate a new document file (DOCX, XLSX, or PDF). "
            "Provide a format and content structure. "
            "Use when: the user asks to create a new report, spreadsheet, "
            "or document from scratch."
        ),
        risk_level=RiskLevel.LOW_WRITE,
        system_id=BUILTIN_SYSTEM_ID,
        method=HttpMethod.POST.value,
        path="/__builtin__/documents/create",
        parameters={
            "format": {
                "type": "string",
                "description": "Output format: docx, xlsx, or pdf",
                "required": True,
            },
            "title": {
                "type": "string",
                "description": "Document title",
                "required": True,
            },
            "content": {
                "type": "object",
                "description": (
                    "Document content. For docx: {paragraphs: [str]}. "
                    "For xlsx: {sheets: [{name: str, rows: [[value]]}]}. "
                    "For pdf: {paragraphs: [str]}."
                ),
                "required": True,
            },
        },
    )


def document_modify_tool() -> ToolDefinition:
    """Tool definition for modifying existing documents."""
    return ToolDefinition(
        endpoint_id="__builtin__document_modify",
        name="document_modify",
        description=(
            "Edit an existing document: append paragraphs, replace text, "
            "update spreadsheet cells, or merge PDFs. "
            "Use when: the user asks to edit, update, or merge documents."
        ),
        risk_level=RiskLevel.LOW_WRITE,
        system_id=BUILTIN_SYSTEM_ID,
        method=HttpMethod.POST.value,
        path="/__builtin__/documents/modify",
        parameters={
            "file_path": {
                "type": "string",
                "description": "Storage path of the document to modify",
                "required": True,
            },
            "operation": {
                "type": "string",
                "description": "Operation: append, replace, update_cells, or merge",
                "required": True,
            },
            "data": {
                "type": "object",
                "description": (
                    "Operation data. append: {paragraphs: [str]}. "
                    "replace: {old_text: str, new_text: str}. "
                    "update_cells: {updates: [{sheet: str, cell: str, value: any}]}. "
                    "merge: {file_paths: [str]}."
                ),
                "required": True,
            },
        },
    )


def document_convert_tool() -> ToolDefinition:
    """Tool definition for converting documents between formats."""
    return ToolDefinition(
        endpoint_id="__builtin__document_convert",
        name="document_convert",
        description=(
            "Convert a document between formats: XLSX to CSV, DOCX to plain text. "
            "Use when: the user needs a document in a different format."
        ),
        risk_level=RiskLevel.READ,
        system_id=BUILTIN_SYSTEM_ID,
        method=HttpMethod.POST.value,
        path="/__builtin__/documents/convert",
        parameters={
            "file_path": {
                "type": "string",
                "description": "Storage path of the source document",
                "required": True,
            },
            "target_format": {
                "type": "string",
                "description": "Target format: csv or text",
                "required": True,
            },
        },
    )


# ---------------------------------------------------------------------------
# Executor
# ---------------------------------------------------------------------------


class DocumentToolExecutor:
    """Execute document tools against the file storage backend.

    Each handler lazily imports the required document library so that
    the dependency is only loaded when the tool is actually invoked.
    """

    _TOOL_NAMES = frozenset({
        "document_read",
        "document_create",
        "document_modify",
        "document_convert",
    })

    def __init__(
        self,
        storage: FileStorageProvider,
        file_repo: FileUploadRepository | None = None,
        settings_repo: SettingsRepository | None = None,
    ) -> None:
        self._storage = storage
        self._file_repo = file_repo
        self._settings_repo = settings_repo

    @classmethod
    def is_document_tool(cls, tool_name: str) -> bool:
        """Return True if *tool_name* is handled by this executor."""
        return tool_name in cls._TOOL_NAMES

    @staticmethod
    def _truncate_result(result: dict[str, Any], max_chars: int) -> dict[str, Any]:
        """Truncate the text content of a read result if it exceeds max_chars."""
        text = result.get("text", "")
        if not text or len(text) <= max_chars:
            return result

        # Truncate at a word boundary
        truncated = text[:max_chars]
        last_space = truncated.rfind(" ")
        if last_space > max_chars * 0.8:
            truncated = truncated[:last_space]

        result = {**result}
        result["text"] = truncated
        result["truncated"] = True
        result["chars_returned"] = len(truncated)
        result["chars_total"] = len(text)
        result["hint"] = (
            "Output was truncated. To read more, call document_read again "
            "with a larger max_chars or use page_start/page_end to read "
            "specific page ranges."
        )
        return result

    async def _resolve_file_path(self, file_path: str) -> str | None:
        """Resolve a user-provided path to the actual storage path.

        When a file is not found at the literal *file_path*, this tries to
        find it by looking up the original filename in the file upload repo.
        """
        if self._file_repo is None:
            logger.warning("document_read: cannot resolve '%s' — file_repo not configured", file_path)
            return None
        import os
        # Extract just the filename from the path
        filename = os.path.basename(file_path)
        upload = await self._file_repo.get_by_filename(filename)
        if upload is not None:
            return upload.storage_path
        logger.warning("document_read: filename '%s' not found in file upload repo", filename)
        return None

    async def _register_file_upload(
        self,
        filename: str,
        storage_path: str,
        content_type: str,
        size_bytes: int,
    ) -> str | None:
        """Create a FileUpload record so the file is downloadable via the API.

        Returns the file ID if the record was created, otherwise ``None``.
        """
        if self._file_repo is None:
            return None

        from flydesk.files.models import FileUpload

        file_id = str(uuid.uuid4())
        upload = FileUpload(
            id=file_id,
            user_id="agent",
            filename=filename,
            content_type=content_type,
            file_size=size_bytes,
            storage_path=storage_path,
            storage_backend="local",
        )
        try:
            await self._file_repo.create(upload)
        except Exception:
            logger.warning("Failed to register FileUpload for %s", filename, exc_info=True)
            return None
        return file_id

    async def execute(self, tool_name: str, arguments: dict[str, Any]) -> dict[str, Any]:
        """Dispatch to the appropriate document handler."""
        handlers: dict[str, Any] = {
            "document_read": self._read,
            "document_create": self._create,
            "document_modify": self._modify,
            "document_convert": self._convert,
        }

        handler = handlers.get(tool_name)
        if handler is None:
            return {"error": f"Unknown document tool: {tool_name}"}

        try:
            return await handler(arguments)
        except Exception as exc:
            logger.error("Document tool %s failed: %s", tool_name, exc, exc_info=True)
            return {"error": str(exc)}

    # ---- Read ----

    async def _read(self, arguments: dict[str, Any]) -> dict[str, Any]:
        file_id: str = arguments.get("file_id", "")
        file_path: str = arguments.get("file_path", "")
        logger.info("document_read called: file_id=%r, file_path=%r", file_id, file_path)

        if not file_id and not file_path:
            return {"error": "Either file_id or file_path is required"}

        default_max = _DEFAULT_MAX_CHARS
        if self._settings_repo is not None:
            try:
                rt = await self._settings_repo.get_llm_runtime_settings()
                default_max = rt.document_read_max_chars
            except Exception:
                logger.debug("LLM runtime settings fetch failed; using default.", exc_info=True)
        max_chars = int(arguments.get("max_chars") or default_max)
        page_start = arguments.get("page_start")
        page_end = arguments.get("page_end")

        resolved_path: str | None = None
        data: bytes | None = None

        # Priority 1: Resolve via file_id (unambiguous)
        if file_id and self._file_repo is not None:
            upload = await self._file_repo.get(file_id)
            if upload is not None:
                resolved_path = upload.storage_path
                logger.info("document_read: resolved file_id=%s → %s", file_id, resolved_path)
            else:
                logger.warning("document_read: file_id=%s not found in repo", file_id)

        # Priority 2: Try file_path directly
        if resolved_path is None and file_path:
            try:
                data = await self._storage.retrieve(file_path)
                resolved_path = file_path
            except FileNotFoundError:
                # Priority 3: Resolve filename via repo lookup
                resolved = await self._resolve_file_path(file_path)
                if resolved is not None:
                    resolved_path = resolved
                    logger.info("document_read: resolved filename '%s' → %s", file_path, resolved_path)

        if resolved_path is None:
            identifier = file_id or file_path
            logger.warning("document_read: could not resolve file: %s", identifier)
            return {"error": f"File not found: {identifier}"}

        # Read from resolved path (skip if already loaded in Priority 2)
        if data is None:
            data = await self._storage.retrieve(resolved_path)

        ext = resolved_path.rsplit(".", 1)[-1].lower() if "." in resolved_path else ""

        if ext == "pdf":
            result = self._read_pdf(
                data,
                page_start=int(page_start) if page_start else None,
                page_end=int(page_end) if page_end else None,
            )
        elif ext == "docx":
            result = self._read_docx(data)
        elif ext == "xlsx":
            result = self._read_xlsx(data)
        elif ext == "pptx":
            result = self._read_pptx(data)
        else:
            return {"error": f"Unsupported file format: .{ext}"}

        return self._truncate_result(result, max_chars)

    @staticmethod
    def _read_pdf(
        data: bytes,
        page_start: int | None = None,
        page_end: int | None = None,
    ) -> dict[str, Any]:
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(data))
        total_pages = len(reader.pages)

        # Convert to 0-based indices; default to all pages
        start = max((page_start or 1) - 1, 0)
        end = min((page_end or total_pages), total_pages)

        pages = [reader.pages[i].extract_text() or "" for i in range(start, end)]
        result: dict[str, Any] = {
            "format": "pdf",
            "total_pages": total_pages,
            "pages_returned": f"{start + 1}-{end}",
            "text": "\n\n".join(pages),
        }
        if end < total_pages:
            result["has_more"] = True
            result["next_page_start"] = end + 1
        return result

    @staticmethod
    def _read_docx(data: bytes) -> dict[str, Any]:
        from docx import Document

        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs]
        tables: list[list[list[str]]] = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text for cell in row.cells])
            tables.append(rows)
        return {
            "format": "docx",
            "paragraphs": paragraphs,
            "tables": tables,
        }

    @staticmethod
    def _read_xlsx(data: bytes) -> dict[str, Any]:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets: dict[str, list[list[Any]]] = {}
        for name in wb.sheetnames:
            ws = wb[name]
            rows: list[list[Any]] = []
            for row in ws.iter_rows(values_only=True):
                rows.append([cell for cell in row])
            sheets[name] = rows
        wb.close()
        return {"format": "xlsx", "sheets": sheets}

    @staticmethod
    def _read_pptx(data: bytes) -> dict[str, Any]:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides: list[dict[str, Any]] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            slides.append({"texts": texts})
        return {"format": "pptx", "slide_count": len(slides), "slides": slides}

    # ---- Create ----

    async def _create(self, arguments: dict[str, Any]) -> dict[str, Any]:
        fmt: str = arguments.get("format", "")
        title: str = arguments.get("title", "")
        content: dict[str, Any] = _coerce_dict(arguments.get("content", {}))

        if not fmt:
            return {"error": "format is required"}
        if not title:
            return {"error": "title is required"}

        if fmt == "docx":
            return await self._create_docx(title, content)
        if fmt == "xlsx":
            return await self._create_xlsx(title, content)
        if fmt == "pptx":
            return await self._create_pptx(title, content)
        if fmt == "pdf":
            return await self._create_pdf(title, content)
        return {"error": f"Unsupported create format: {fmt}"}

    async def _create_docx(self, title: str, content: dict[str, Any]) -> dict[str, Any]:
        from docx import Document

        doc = Document()
        doc.add_heading(title, level=0)
        for para in content.get("paragraphs", []):
            doc.add_paragraph(str(para))

        buf = io.BytesIO()
        doc.save(buf)
        file_bytes = buf.getvalue()
        filename = f"{title}.docx"
        path = await self._storage.store(filename, file_bytes, _CONTENT_TYPES["docx"])
        result: dict[str, Any] = {"format": "docx", "file_path": path, "filename": filename}
        file_id = await self._register_file_upload(filename, path, _CONTENT_TYPES["docx"], len(file_bytes))
        if file_id:
            result["file_id"] = file_id
        return result

    async def _create_xlsx(self, title: str, content: dict[str, Any]) -> dict[str, Any]:
        from openpyxl import Workbook

        wb = Workbook()
        # Remove the default sheet so we start clean
        wb.remove(wb.active)  # type: ignore[arg-type]

        sheets_data = content.get("sheets", [])
        if not sheets_data:
            # Fallback: create a single sheet named after the title
            ws = wb.create_sheet(title=title)
            for row_data in content.get("rows", []):
                ws.append(row_data)
        else:
            for sheet_info in sheets_data:
                ws = wb.create_sheet(title=sheet_info.get("name", "Sheet"))
                for row_data in sheet_info.get("rows", []):
                    ws.append(row_data)

        buf = io.BytesIO()
        wb.save(buf)
        file_bytes = buf.getvalue()
        filename = f"{title}.xlsx"
        path = await self._storage.store(filename, file_bytes, _CONTENT_TYPES["xlsx"])
        result: dict[str, Any] = {"format": "xlsx", "file_path": path, "filename": filename}
        file_id = await self._register_file_upload(filename, path, _CONTENT_TYPES["xlsx"], len(file_bytes))
        if file_id:
            result["file_id"] = file_id
        return result

    async def _create_pdf(self, title: str, content: dict[str, Any]) -> dict[str, Any]:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=letter)
        styles = getSampleStyleSheet()

        story: list[Any] = []
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 12))
        for para in content.get("paragraphs", []):
            story.append(Paragraph(str(para), styles["BodyText"]))
            story.append(Spacer(1, 6))

        doc.build(story)
        file_bytes = buf.getvalue()
        filename = f"{title}.pdf"
        path = await self._storage.store(filename, file_bytes, _CONTENT_TYPES["pdf"])
        result: dict[str, Any] = {"format": "pdf", "file_path": path, "filename": filename}
        file_id = await self._register_file_upload(filename, path, _CONTENT_TYPES["pdf"], len(file_bytes))
        if file_id:
            result["file_id"] = file_id
        return result

    async def _create_pptx(self, title: str, content: dict[str, Any]) -> dict[str, Any]:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = content.get("subtitle", "")

        # Content slides
        body_layout = prs.slide_layouts[1]
        for slide_data in content.get("slides", []):
            s = prs.slides.add_slide(body_layout)
            s.shapes.title.text = slide_data.get("title", "")
            body = s.placeholders[1]
            tf = body.text_frame
            for i, bullet in enumerate(slide_data.get("bullets", [])):
                if i == 0:
                    tf.text = str(bullet)
                else:
                    tf.add_paragraph().text = str(bullet)

        buf = io.BytesIO()
        prs.save(buf)
        file_bytes = buf.getvalue()
        filename = f"{title}.pptx"
        path = await self._storage.store(filename, file_bytes, _CONTENT_TYPES["pptx"])
        result: dict[str, Any] = {"format": "pptx", "file_path": path, "filename": filename}
        file_id = await self._register_file_upload(filename, path, _CONTENT_TYPES["pptx"], len(file_bytes))
        if file_id:
            result["file_id"] = file_id
        return result

    # ---- Modify ----

    async def _modify(self, arguments: dict[str, Any]) -> dict[str, Any]:
        file_path: str = arguments.get("file_path", "")
        operation: str = arguments.get("operation", "")
        data: dict[str, Any] = _coerce_dict(arguments.get("data", {}))

        if not file_path:
            return {"error": "file_path is required"}
        if not operation:
            return {"error": "operation is required"}

        ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""

        if operation == "append" and ext == "docx":
            return await self._modify_docx_append(file_path, data)
        if operation == "replace" and ext == "docx":
            return await self._modify_docx_replace(file_path, data)
        if operation == "update_cells" and ext == "xlsx":
            return await self._modify_xlsx_update_cells(file_path, data)
        if operation == "merge" and ext == "pdf":
            return await self._modify_pdf_merge(file_path, data)

        return {"error": f"Unsupported operation '{operation}' for .{ext} files"}

    async def _modify_docx_append(
        self, file_path: str, data: dict[str, Any]
    ) -> dict[str, Any]:
        from docx import Document

        raw = await self._storage.retrieve(file_path)
        doc = Document(io.BytesIO(raw))
        for para in data.get("paragraphs", []):
            doc.add_paragraph(str(para))

        buf = io.BytesIO()
        doc.save(buf)
        new_path = await self._storage.store(
            file_path.rsplit("/", 1)[-1], buf.getvalue(), _CONTENT_TYPES["docx"]
        )
        return {"file_path": new_path, "operation": "append", "status": "ok"}

    async def _modify_docx_replace(
        self, file_path: str, data: dict[str, Any]
    ) -> dict[str, Any]:
        from docx import Document

        old_text: str = data.get("old_text", "")
        new_text: str = data.get("new_text", "")
        if not old_text:
            return {"error": "old_text is required for replace operation"}

        raw = await self._storage.retrieve(file_path)
        doc = Document(io.BytesIO(raw))

        count = 0
        for para in doc.paragraphs:
            if old_text in para.text:
                for run in para.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        count += 1

        buf = io.BytesIO()
        doc.save(buf)
        new_path = await self._storage.store(
            file_path.rsplit("/", 1)[-1], buf.getvalue(), _CONTENT_TYPES["docx"]
        )
        return {
            "file_path": new_path,
            "operation": "replace",
            "replacements": count,
            "status": "ok",
        }

    async def _modify_xlsx_update_cells(
        self, file_path: str, data: dict[str, Any]
    ) -> dict[str, Any]:
        from openpyxl import load_workbook

        raw = await self._storage.retrieve(file_path)
        wb = load_workbook(io.BytesIO(raw))

        updates = data.get("updates", [])
        for upd in updates:
            sheet_name = upd.get("sheet", wb.sheetnames[0])
            cell = upd.get("cell", "A1")
            value = upd.get("value")
            ws = wb[sheet_name]
            ws[cell] = value

        buf = io.BytesIO()
        wb.save(buf)
        wb.close()
        new_path = await self._storage.store(
            file_path.rsplit("/", 1)[-1], buf.getvalue(), _CONTENT_TYPES["xlsx"]
        )
        return {
            "file_path": new_path,
            "operation": "update_cells",
            "updates_applied": len(updates),
            "status": "ok",
        }

    async def _modify_pdf_merge(
        self, file_path: str, data: dict[str, Any]
    ) -> dict[str, Any]:
        from PyPDF2 import PdfMerger

        merger = PdfMerger()
        # Append the primary file first
        primary = await self._storage.retrieve(file_path)
        merger.append(io.BytesIO(primary))

        extra_paths: list[str] = data.get("file_paths", [])
        for extra_path in extra_paths:
            extra_data = await self._storage.retrieve(extra_path)
            merger.append(io.BytesIO(extra_data))

        buf = io.BytesIO()
        merger.write(buf)
        merger.close()

        new_path = await self._storage.store(
            file_path.rsplit("/", 1)[-1], buf.getvalue(), _CONTENT_TYPES["pdf"]
        )
        return {
            "file_path": new_path,
            "operation": "merge",
            "merged_count": 1 + len(extra_paths),
            "status": "ok",
        }

    # ---- Convert ----

    async def _convert(self, arguments: dict[str, Any]) -> dict[str, Any]:
        file_path: str = arguments.get("file_path", "")
        target_format: str = arguments.get("target_format", "")

        if not file_path:
            return {"error": "file_path is required"}
        if not target_format:
            return {"error": "target_format is required"}

        ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""

        if ext == "xlsx" and target_format == "csv":
            return await self._convert_xlsx_to_csv(file_path)
        if ext == "docx" and target_format == "text":
            return await self._convert_docx_to_text(file_path)

        return {"error": f"Unsupported conversion: .{ext} -> {target_format}"}

    async def _convert_xlsx_to_csv(self, file_path: str) -> dict[str, Any]:
        from openpyxl import load_workbook

        raw = await self._storage.retrieve(file_path)
        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        ws = wb[wb.sheetnames[0]]

        buf = io.StringIO()
        writer = csv.writer(buf)
        for row in ws.iter_rows(values_only=True):
            writer.writerow([cell if cell is not None else "" for cell in row])
        wb.close()

        csv_bytes = buf.getvalue().encode("utf-8")
        base_name = file_path.rsplit("/", 1)[-1].rsplit(".", 1)[0]
        filename = f"{base_name}.csv"
        new_path = await self._storage.store(filename, csv_bytes, "text/csv")
        return {"file_path": new_path, "target_format": "csv", "filename": filename}

    async def _convert_docx_to_text(self, file_path: str) -> dict[str, Any]:
        from docx import Document

        raw = await self._storage.retrieve(file_path)
        doc = Document(io.BytesIO(raw))
        text = "\n".join(p.text for p in doc.paragraphs)

        text_bytes = text.encode("utf-8")
        base_name = file_path.rsplit("/", 1)[-1].rsplit(".", 1)[0]
        filename = f"{base_name}.txt"
        new_path = await self._storage.store(filename, text_bytes, "text/plain")
        return {"file_path": new_path, "target_format": "text", "filename": filename}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_CONTENT_TYPES: dict[str, str] = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
}
