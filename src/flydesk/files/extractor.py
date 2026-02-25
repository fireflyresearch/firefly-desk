# Copyright 2026 Firefly Software Solutions Inc
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0

"""Content extraction from uploaded files."""

from __future__ import annotations

import io
import logging

logger = logging.getLogger(__name__)

# Content-type constants
_PDF_TYPES = frozenset({
    "application/pdf",
})
_DOCX_TYPES = frozenset({
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
})
_XLSX_TYPES = frozenset({
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
})
_PPTX_TYPES = frozenset({
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
})


class ContentExtractor:
    """Extract text content from uploaded files."""

    async def extract(
        self, filename: str, content: bytes, content_type: str
    ) -> str | None:
        """Extract text content from a file.

        Returns the extracted text for supported types, or ``None`` for
        binary formats that are not yet supported.
        """
        if content_type.startswith("text/"):
            return content.decode("utf-8", errors="replace")
        if content_type == "application/json":
            return content.decode("utf-8", errors="replace")
        if content_type in _PDF_TYPES:
            return self._extract_pdf(content)
        if content_type in _DOCX_TYPES:
            return self._extract_docx(content)
        if content_type in _XLSX_TYPES:
            return self._extract_xlsx(content)
        if content_type in _PPTX_TYPES:
            return self._extract_pptx(content)
        return None

    @staticmethod
    def _extract_pdf(content: bytes) -> str | None:
        """Extract text from a PDF document using PyPDF2."""
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            pages: list[str] = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages) if pages else None
        except Exception:
            logger.warning("Failed to extract text from PDF", exc_info=True)
            return None

    @staticmethod
    def _extract_docx(content: bytes) -> str | None:
        """Extract paragraph text from a DOCX document using python-docx."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs) if paragraphs else None
        except Exception:
            logger.warning("Failed to extract text from DOCX", exc_info=True)
            return None

    @staticmethod
    def _extract_xlsx(content: bytes) -> str | None:
        """Extract cell values from an XLSX workbook using openpyxl."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            sheets: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[str] = []
                rows.append(f"--- {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    rows.append("\t".join(cells))
                sheets.append("\n".join(rows))
            wb.close()
            return "\n\n".join(sheets) if sheets else None
        except Exception:
            logger.warning("Failed to extract text from XLSX", exc_info=True)
            return None

    @staticmethod
    def _extract_pptx(content: bytes) -> str | None:
        """Extract text from a PPTX presentation using python-pptx."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            slides: list[str] = []
            for slide in prs.slides:
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    slides.append("\n".join(texts))
            return "\n\n".join(slides) if slides else None
        except Exception:
            logger.warning("Failed to extract text from PPTX", exc_info=True)
            return None
