# Copyright 2026 Firefly Software Solutions Inc
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0

"""Tests for ContentExtractor."""

from __future__ import annotations

import io
import json

import pytest

from flydesk.files.extractor import ContentExtractor


@pytest.fixture
def extractor():
    return ContentExtractor()


class TestContentExtractorExisting:
    """Tests for existing plain-text and JSON extraction."""

    async def test_plain_text(self, extractor):
        """Plain text files are decoded and returned."""
        content = b"Hello, world!"
        result = await extractor.extract("readme.txt", content, "text/plain")
        assert result == "Hello, world!"

    async def test_html_text(self, extractor):
        """HTML content (text/html) is decoded and returned."""
        content = b"<h1>Title</h1>"
        result = await extractor.extract("page.html", content, "text/html")
        assert result == "<h1>Title</h1>"

    async def test_json(self, extractor):
        """JSON content is decoded and returned as a string."""
        data = {"key": "value", "number": 42}
        content = json.dumps(data).encode()
        result = await extractor.extract("data.json", content, "application/json")
        assert result is not None
        assert "key" in result
        assert "value" in result

    async def test_unknown_binary_returns_none(self, extractor):
        """Unknown binary types return None."""
        result = await extractor.extract(
            "image.png", b"\x89PNG\r\n", "image/png"
        )
        assert result is None


class TestContentExtractorPDF:
    """Tests for PDF extraction via PyPDF2."""

    async def test_extract_pdf(self, extractor):
        """PDF files are extracted with text content preserved."""
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=letter)
        c.drawString(72, 700, "Firefly PDF Test Content")
        c.drawString(72, 680, "Second line of text")
        c.showPage()
        c.save()
        pdf_bytes = buf.getvalue()

        result = await extractor.extract(
            "report.pdf", pdf_bytes, "application/pdf"
        )
        assert result is not None
        assert "Firefly PDF Test Content" in result
        assert "Second line" in result

    async def test_extract_pdf_multipage(self, extractor):
        """Multi-page PDFs have text extracted from all pages."""
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=letter)
        c.drawString(72, 700, "Page one content")
        c.showPage()
        c.drawString(72, 700, "Page two content")
        c.showPage()
        c.save()
        pdf_bytes = buf.getvalue()

        result = await extractor.extract(
            "multi.pdf", pdf_bytes, "application/pdf"
        )
        assert result is not None
        assert "Page one" in result
        assert "Page two" in result


class TestContentExtractorDOCX:
    """Tests for DOCX extraction via python-docx."""

    async def test_extract_docx(self, extractor):
        """DOCX files have paragraph text extracted."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Firefly DOCX Test Content")
        doc.add_paragraph("Another paragraph here")
        buf = io.BytesIO()
        doc.save(buf)
        docx_bytes = buf.getvalue()

        result = await extractor.extract(
            "document.docx",
            docx_bytes,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        assert result is not None
        assert "Firefly DOCX Test Content" in result
        assert "Another paragraph" in result

    async def test_extract_docx_legacy_content_type(self, extractor):
        """Legacy application/msword content type is also handled."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Legacy format test")
        buf = io.BytesIO()
        doc.save(buf)
        docx_bytes = buf.getvalue()

        result = await extractor.extract(
            "legacy.doc", docx_bytes, "application/msword"
        )
        assert result is not None
        assert "Legacy format test" in result


class TestContentExtractorXLSX:
    """Tests for XLSX extraction via openpyxl."""

    async def test_extract_xlsx(self, extractor):
        """XLSX files have cell values extracted."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sales"
        ws["A1"] = "Product"
        ws["B1"] = "Revenue"
        ws["A2"] = "Widget"
        ws["B2"] = 1500
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        result = await extractor.extract(
            "data.xlsx",
            xlsx_bytes,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        assert result is not None
        assert "Product" in result
        assert "Revenue" in result
        assert "Widget" in result
        assert "1500" in result

    async def test_extract_xlsx_multiple_sheets(self, extractor):
        """Multi-sheet XLSX workbooks have all sheets extracted."""
        from openpyxl import Workbook

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1["A1"] = "Alpha"
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Beta"
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        result = await extractor.extract(
            "multi.xlsx",
            xlsx_bytes,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        assert result is not None
        assert "Alpha" in result
        assert "Beta" in result
        assert "Sheet1" in result
        assert "Sheet2" in result

    async def test_extract_xlsx_legacy_content_type(self, extractor):
        """Legacy application/vnd.ms-excel content type is also handled."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Legacy Excel"
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        result = await extractor.extract(
            "legacy.xls", xlsx_bytes, "application/vnd.ms-excel"
        )
        assert result is not None
        assert "Legacy Excel" in result


class TestContentExtractorPPTX:
    """Tests for PPTX extraction via python-pptx."""

    async def test_extract_pptx(self, extractor):
        """PPTX files have slide text extracted."""
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Firefly PPTX Title"
        slide.placeholders[1].text = "Slide body content here"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        result = await extractor.extract(
            "presentation.pptx",
            pptx_bytes,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        assert result is not None
        assert "Firefly PPTX Title" in result
        assert "Slide body content" in result

    async def test_extract_pptx_multiple_slides(self, extractor):
        """Multi-slide PPTX presentations have all slides extracted."""
        from pptx import Presentation

        prs = Presentation()
        for i in range(1, 4):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Slide {i} title"
            slide.placeholders[1].text = f"Content for slide {i}"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        result = await extractor.extract(
            "multi.pptx",
            pptx_bytes,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        assert result is not None
        assert "Slide 1 title" in result
        assert "Slide 3 title" in result
        assert "Content for slide 2" in result

    async def test_extract_pptx_legacy_content_type(self, extractor):
        """Legacy application/vnd.ms-powerpoint content type is also handled."""
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Legacy PowerPoint"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        result = await extractor.extract(
            "legacy.ppt", pptx_bytes, "application/vnd.ms-powerpoint"
        )
        assert result is not None
        assert "Legacy PowerPoint" in result
